# -*- coding: utf-8 -*-
"""
基础文档解析组件 - 核心预装依赖
包含PDF、Word、Excel、PPT的基础解析功能
"""

import asyncio
from pathlib import Path
from typing import Dict, Any, List, Optional
import tempfile
import os

# 可选依赖 - 延迟导入，允许模块在缺少部分库时仍能加载
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import xlrd
except ImportError:
    xlrd = None

try:
    import pptx
except ImportError:
    pptx = None

from ...utils.logging import logger
from . import LocalComponent


class BasicParserComponent(LocalComponent):
    """
    基础文档解析组件
    支持原生PDF、Word、Excel、PPT文本提取
    """
    
    def __init__(self):
        super().__init__(
            component_id="basic_parser",
            name="基础文档解析",
            description="原生PDF/Word/Excel/PPT文本提取（核心预装）",
            install_size_mb=0.0  # 预装，无额外空间
        )
        
        # 预装组件无需安装依赖
        self.required_packages = []
        self.required_system_deps = []
        
        # 支持的 MIME 类型
        self.supported_mimes = [
            "application/pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ]
        
        # 文件扩展名映射
        self.extension_mimes = {
            ".pdf": "application/pdf",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ".xls": "application/vnd.ms-excel",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }
    
    async def initialize(self, manager) -> bool:
        """初始化组件"""
        try:
            self.is_installed = True
            self.is_enabled = True
            
            # 验证基础库可用性
            missing = []
            if fitz is None:
                missing.append("PyMuPDF (fitz)")
            if docx is None:
                missing.append("python-docx")
            if openpyxl is None:
                missing.append("openpyxl")
            if pptx is None:
                missing.append("python-pptx")
            
            if missing:
                from ...utils.logging import logger
                logger.warning(
                    f"基础解析组件部分依赖缺失: {', '.join(missing)}。"
                    f"支持的功能将受限。"
                )
                # 仍然标记为已安装，只是功能受限
                return True
            
            await self._verify_dependencies()
            
            logger.info("基础文档解析组件初始化成功")
            return True
            
        except Exception as e:
            logger.error(f"基础解析组件初始化失败: {e}")
            return False
    
    async def _verify_dependencies(self):
        """验证依赖库"""
        try:
            # 简单测试各个库是否可用
            if fitz is not None:
                doc = fitz.open()
                doc.close()
                
                with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                    doc = fitz.open()
                    page = doc.new_page()
                    doc.save(tmp.name)
                    doc.close()
                    os.unlink(tmp.name)
            
            logger.debug("基础解析依赖库验证通过")
            
        except Exception as e:
            logger.error(f"依赖库验证失败: {e}")
            raise
    
    async def parse_document(
        self, 
        file_path: str, 
        options: Dict[str, Any] = None
    ) -> 'ParseResult':
        """解析文档"""
        options = options or {}
        file_path_obj = Path(file_path)
        
        if not file_path_obj.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        # 根据文件扩展名选择解析器
        extension = file_path_obj.suffix.lower()
        mime_type = self.extension_mimes.get(extension)
        
        if not mime_type:
            raise ValueError(f"不支持的文件格式: {extension}")
        
        # 执行解析
        if mime_type == "application/pdf":
            return await self._parse_pdf(file_path, options)
        elif mime_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document"]:
            return await self._parse_docx(file_path, options)
        elif mime_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel"]:
            return await self._parse_excel(file_path, options)
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return await self._parse_pptx(file_path, options)
        else:
            raise ValueError(f"未实现的MIME类型: {mime_type}")
    
    async def _parse_pdf(self, file_path: str, options: Dict[str, Any]) -> 'ParseResult':
        """解析PDF文档"""
        text_parts = []
        tables = []
        
        try:
            # 使用PyMuPDF提取文本
            def extract_with_fitz():
                text_parts_fitz = []
                with fitz.open(file_path) as doc:
                    for page_num, page in enumerate(doc, 1):
                        text = page.get_text()
                        if text.strip():
                            text_parts_fitz.append(f"=== 第{page_num}页 ===\n{text}")
                return text_parts_fitz
            
            # 使用pdfplumber提取表格
            def extract_tables():
                tables_data = []
                try:
                    with pdfplumber.open(file_path) as pdf:
                        for page_num, page in enumerate(pdf.pages, 1):
                            page_tables = page.extract_tables()
                            for table_idx, table in enumerate(page_tables):
                                if table and len(table) > 1:  # 至少有两行
                                    tables_data.append({
                                        "page": page_num,
                                        "table_id": f"table_{page_num}_{table_idx}",
                                        "data": table
                                    })
                except Exception as e:
                    logger.debug(f"提取表格失败: {e}")
                return tables_data
            
            # 并行执行
            fitz_task = asyncio.create_task(asyncio.to_thread(extract_with_fitz))
            tables_task = asyncio.create_task(asyncio.to_thread(extract_tables))
            
            text_parts = await fitz_task
            tables = await tables_task
            
            # 完整文本
            full_text = "\n".join(text_parts)
            
            # 生成Markdown格式
            markdown = self._text_to_markdown(full_text, tables)
            
            return ParseResult(
                text=full_text,
                markdown=markdown,
                tables=tables,
                engine_info={
                    "engine": "basic_pdf_fitz_plumber",
                    "page_count": len(text_parts),
                    "table_count": len(tables),
                    "is_cloud": False,
                    "cost": 0.0
                },
                metadata={
                    "file_type": "pdf",
                    "parser": "fitz_plumber"
                }
            )
            
        except Exception as e:
            logger.error(f"PDF解析失败: {e}")
            raise
    
    async def _parse_docx(self, file_path: str, options: Dict[str, Any]) -> 'ParseResult':
        """解析Word文档"""
        try:
            def extract_docx():
                doc = docx.Document(file_path)
                text_parts = []
                tables = []
                
                # 提取段落
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
                
                # 提取表格
                for table_idx, table in enumerate(doc.tables):
                    table_data = []
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text)
                        table_data.append(row_data)
                    
                    if table_data:
                        tables.append({
                            "table_id": f"table_{table_idx}",
                            "data": table_data
                        })
                
                return text_parts, tables
            
            text_parts, tables = await asyncio.to_thread(extract_docx)
            full_text = "\n".join(text_parts)
            markdown = self._text_to_markdown(full_text, tables)
            
            return ParseResult(
                text=full_text,
                markdown=markdown,
                tables=tables,
                engine_info={
                    "engine": "basic_docx_python_docx",
                    "paragraph_count": len(text_parts),
                    "table_count": len(tables),
                    "is_cloud": False,
                    "cost": 0.0
                },
                metadata={
                    "file_type": "docx",
                    "parser": "python-docx"
                }
            )
            
        except Exception as e:
            logger.error(f"Word解析失败: {e}")
            raise
    
    async def _parse_excel(self, file_path: str, options: Dict[str, Any]) -> 'ParseResult':
        """解析Excel文档"""
        try:
            file_path_obj = Path(file_path)
            
            if file_path_obj.suffix.lower() == '.xlsx':
                # 新格式
                def extract_xlsx():
                    wb = openpyxl.load_workbook(file_path)
                    text_parts = []
                    tables = []
                    
                    for sheet_name in wb.sheetnames:
                        sheet = wb[sheet_name]
                        text_parts.append(f"=== 工作表: {sheet_name} ===")
                        
                        # 提取数据
                        table_data = []
                        for row in sheet.iter_rows(values_only=True):
                            # 过滤空行
                            if any(cell is not None and str(cell).strip() for cell in row):
                                table_data.append([str(cell) if cell is not None else "" for cell in row])
                        
                        if table_data:
                            tables.append({
                                "table_id": f"sheet_{sheet_name}",
                                "sheet_name": sheet_name,
                                "data": table_data
                            })
                            text_parts.append(f"工作表包含 {len(table_data)} 行数据")
                    
                    return text_parts, tables
                
                text_parts, tables = await asyncio.to_thread(extract_xlsx)
                
            else:
                # 旧格式.xls
                def extract_xls():
                    wb = xlrd.open_workbook(file_path)
                    text_parts = []
                    tables = []
                    
                    for sheet_idx in range(wb.nsheets):
                        sheet = wb.sheet_by_index(sheet_idx)
                        sheet_name = wb.sheet_names()[sheet_idx]
                        
                        text_parts.append(f"=== 工作表: {sheet_name} ===")
                        
                        # 提取数据
                        table_data = []
                        for row_idx in range(sheet.nrows):
                            row_data = []
                            for col_idx in range(sheet.ncols):
                                cell = sheet.cell(row_idx, col_idx)
                                row_data.append(str(cell.value) if cell.value else "")
                            
                            # 过滤空行
                            if any(cell.strip() for cell in row_data):
                                table_data.append(row_data)
                        
                        if table_data:
                            tables.append({
                                "table_id": f"sheet_{sheet_name}",
                                "sheet_name": sheet_name,
                                "data": table_data
                            })
                            text_parts.append(f"工作表包含 {len(table_data)} 行数据")
                    
                    return text_parts, tables
                
                text_parts, tables = await asyncio.to_thread(extract_xls)
            
            full_text = "\n".join(text_parts)
            markdown = self._text_to_markdown(full_text, tables)
            
            return ParseResult(
                text=full_text,
                markdown=markdown,
                tables=tables,
                engine_info={
                    "engine": "basic_excel_openpyxl_xlrd",
                    "sheet_count": len(tables),
                    "is_cloud": False,
                    "cost": 0.0
                },
                metadata={
                    "file_type": "excel",
                    "parser": "openpyxl_xlrd"
                }
            )
            
        except Exception as e:
            logger.error(f"Excel解析失败: {e}")
            raise
    
    async def _parse_pptx(self, file_path: str, options: Dict[str, Any]) -> 'ParseResult':
        """解析PowerPoint文档"""
        try:
            def extract_pptx():
                prs = pptx.Presentation(file_path)
                text_parts = []
                
                for slide_idx, slide in enumerate(prs.slides, 1):
                    text_parts.append(f"=== 第{slide_idx}页幻灯片 ===")
                    slide_text = []
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    if slide_text:
                        text_parts.extend(slide_text)
                    else:
                        text_parts.append("(空白页面)")
                
                return text_parts
            
            text_parts = await asyncio.to_thread(extract_pptx)
            full_text = "\n".join(text_parts)
            markdown = self._text_to_markdown(full_text, [])
            
            return ParseResult(
                text=full_text,
                markdown=markdown,
                tables=[],
                engine_info={
                    "engine": "basic_pptx_python_pptx",
                    "slide_count": len(text_parts),
                    "is_cloud": False,
                    "cost": 0.0
                },
                metadata={
                    "file_type": "pptx",
                    "parser": "python-pptx"
                }
            )
            
        except Exception as e:
            logger.error(f"PowerPoint解析失败: {e}")
            raise
    
    def _text_to_markdown(self, text: str, tables: List[Dict]) -> str:
        """将文本和表格转换为Markdown格式"""
        lines = ["# 文档内容\n"]
        lines.append(text)
        lines.append("\n")
        
        # 添加表格
        for i, table_data in enumerate(tables):
            lines.append(f"\n## 表格 {i+1}\n")
            if "table_id" in table_data:
                lines.append(f"**表格ID**: `{table_data['table_id']}`\n")
            
            table = table_data.get("data", [])
            if table:
                # 表格Markdown
                for j, row in enumerate(table):
                    line = "| " + " | ".join(str(cell) for cell in row) + " |"
                    lines.append(line)
                    
                    # 第一行后添加分隔线
                    if j == 0:
                        sep = "| " + " | ".join("---" for _ in row) + " |"
                        lines.append(sep)
                
                lines.append("")
        
        return "\n".join(lines)
    
    async def check_dependencies(self) -> bool:
        """检查依赖 - 核心组件无需额外依赖"""
        return True
    
    async def get_capabilities(self) -> Dict[str, Any]:
        """获取组件能力"""
        return {
            "supported_formats": [
                "pdf", "docx", "xlsx", "xls", "pptx"
            ],
            "features": [
                "text_extraction",
                "table_extraction", 
                "basic_formatting"
            ],
            "limitations": [
                "不支持扫描件OCR",
                "不支持复杂版面分析",
                "不支持公式识别"
            ],
            "performance": {
                "speed": "fast",
                "accuracy": "medium",
                "memory_usage": "low"
            }
        }